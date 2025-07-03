import os
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

class PPTGenerator:
    def __init__(self, folder_path, output_filename="generated_presentation.pptx"):
        """
        初始化PPT生成器
        
        Args:
            folder_path (str): 包含图片和txt文件的文件夹路径
            output_filename (str): 输出的PPT文件名
        """
        self.folder_path = folder_path
        self.output_filename = output_filename
        self.presentation = Presentation()
        self.slide_width = self.presentation.slide_width
        self.slide_height = self.presentation.slide_height
        
        # 支持的图片格式
        self.image_extensions = ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']
        
    def get_image_files(self):
        """获取文件夹中所有的图片文件"""
        image_files = []
        for ext in self.image_extensions:
            pattern = os.path.join(self.folder_path, f"*{ext}")
            image_files.extend(glob.glob(pattern, recursive=False))
            # 同时查找大写扩展名
            pattern = os.path.join(self.folder_path, f"*{ext.upper()}")
            image_files.extend(glob.glob(pattern, recursive=False))
        
        return sorted(image_files)
    
    def get_txt_files(self):
        """获取文件夹中所有的txt文件"""
        txt_pattern = os.path.join(self.folder_path, "*.txt")
        return glob.glob(txt_pattern)
    
    def read_overall_view_info(self):
        """读取txt文件中的overall view信息"""
        txt_files = self.get_txt_files()
        overall_info = ""
        
        for txt_file in txt_files:
            try:
                with open(txt_file, 'r', encoding='utf-8') as f:
                    content = f.read()
                    overall_info += f"=== {os.path.basename(txt_file)} ===\n"
                    overall_info += content + "\n\n"
            except UnicodeDecodeError:
                # 尝试其他编码
                try:
                    with open(txt_file, 'r', encoding='gbk') as f:
                        content = f.read()
                        overall_info += f"=== {os.path.basename(txt_file)} ===\n"
                        overall_info += content + "\n\n"
                except:
                    print(f"无法读取文件: {txt_file}")
        
        return overall_info if overall_info else "未找到相关信息文件"
    
    def create_title_slide(self):
        """创建标题页"""
        title_slide_layout = self.presentation.slide_layouts[0]  # 标题页布局
        slide = self.presentation.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "图片分析报告"
        subtitle.text = f"基于文件夹: {os.path.basename(self.folder_path)}"
        
        return slide
    
    def create_overall_view_slide(self):
        """创建Overall View页面"""
        # 使用标题和内容布局
        content_slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        title.text = "Overall View"
        
        # 添加内容
        content_placeholder = slide.placeholders[1]
        overall_info = self.read_overall_view_info()
        
        text_frame = content_placeholder.text_frame
        text_frame.text = overall_info
        
        # 设置文本格式
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = 'Arial'
        
        return slide
    
    def create_image_slide(self, image_path):
        """
        为单张图片创建幻灯片
        
        Args:
            image_path (str): 图片文件路径
        """
        # 使用空白布局
        blank_slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(blank_slide_layout)
        
        # 获取图片名称作为标题（不包含扩展名）
        image_name = Path(image_path).stem
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(1)
        )
        
        title_frame = title_shape.text_frame
        title_frame.text = image_name
        
        # 设置标题格式
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.bold = True
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.name = 'Arial'
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # 计算图片区域（从标题下方到页尾）
        image_top = Inches(1.8)  # 标题下方
        image_height = self.slide_height - image_top - Inches(0.3)  # 页尾留0.3英寸边距
        image_width = Inches(8.5)  # 图片宽度
        image_left = (self.slide_width - image_width) / 2  # 居中
        
        try:
            # 添加图片
            picture = slide.shapes.add_picture(
                image_path,
                left=image_left,
                top=image_top,
                width=image_width,
                height=image_height
            )
            
            print(f"已添加图片: {image_name}")
            
        except Exception as e:
            print(f"添加图片失败 {image_path}: {str(e)}")
            
            # 如果图片添加失败，添加错误信息
            error_shape = slide.shapes.add_textbox(
                left=image_left,
                top=image_top,
                width=image_width,
                height=image_height
            )
            
            error_frame = error_shape.text_frame
            error_frame.text = f"无法加载图片: {image_name}\n错误: {str(e)}"
            error_paragraph = error_frame.paragraphs[0]
            error_paragraph.alignment = PP_ALIGN.CENTER
            error_paragraph.font.size = Pt(16)
            error_paragraph.font.color.rgb = RGBColor(255, 0, 0)  # 红色
        
        return slide
    
    def generate_presentation(self):
        """生成完整的演示文稿"""
        print(f"开始处理文件夹: {self.folder_path}")
        
        # 检查文件夹是否存在
        if not os.path.exists(self.folder_path):
            print(f"错误: 文件夹不存在 - {self.folder_path}")
            return False
        
        # 创建标题页
        self.create_title_slide()
        print("已创建标题页")
        
        # 创建Overall View页
        self.create_overall_view_slide()
        print("已创建Overall View页")
        
        # 获取所有图片文件
        image_files = self.get_image_files()
        
        if not image_files:
            print("警告: 未找到任何图片文件")
            return False
        
        print(f"找到 {len(image_files)} 张图片")
        
        # 为每张图片创建幻灯片
        for image_file in image_files:
            self.create_image_slide(image_file)
        
        # 保存演示文稿
        try:
            self.presentation.save(self.output_filename)
            print(f"演示文稿已保存: {self.output_filename}")
            return True
        except Exception as e:
            print(f"保存演示文稿失败: {str(e)}")
            return False
    
    def analyze_iou_scores(self, image_files):
        """
        分析IOU评分（预留功能）
        
        Args:
            image_files (list): 图片文件列表
            
        Returns:
            dict: 按IOU分数分类的图片字典
        """
        # TODO: 实现IOU分析逻辑
        # 这里可以根据文件名或外部数据源来获取IOU分数
        
        categories = {
            'high_iou': [],      # IOU > 0.7
            'medium_iou': [],    # 0.5 < IOU <= 0.7
            'low_iou': []        # IOU <= 0.5
        }
        
        for image_file in image_files:
            # 示例: 从文件名中提取IOU分数
            # 假设文件名格式为: "image_name_iou_0.85.jpg"
            filename = os.path.basename(image_file)
            iou_match = re.search(r'iou[_-](\d+\.?\d*)', filename.lower())
            
            if iou_match:
                iou_score = float(iou_match.group(1))
                if iou_score > 0.7:
                    categories['high_iou'].append(image_file)
                elif iou_score > 0.5:
                    categories['medium_iou'].append(image_file)
                else:
                    categories['low_iou'].append(image_file)
            else:
                # 如果无法从文件名提取IOU，默认放入medium_iou
                categories['medium_iou'].append(image_file)
        
        return categories


def main():
    """主函数示例"""
    # 配置参数
    folder_path = input("请输入包含图片的文件夹路径: ").strip()
    if not folder_path:
        folder_path = "./images"  # 默认路径
    
    output_filename = input("请输入输出PPT文件名 (默认: presentation.pptx): ").strip()
    if not output_filename:
        output_filename = "presentation.pptx"
    
    # 创建PPT生成器
    generator = PPTGenerator(folder_path, output_filename)
    
    # 生成演示文稿
    success = generator.generate_presentation()
    
    if success:
        print("\n✅ PPT生成完成!")
        print(f"📁 文件位置: {os.path.abspath(output_filename)}")
    else:
        print("\n❌ PPT生成失败，请检查错误信息")


if __name__ == "__main__":
    main()